"""
GEN_AI_TOOL project
Router and AI responses comparison tool done with flask

mrbacco04@gmail.com
Q2, 2026

"""

"""Local document indexing and retrieval for grounded prompting."""

import hashlib
import json
import math
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import threading
import uuid
from datetime import datetime, timezone

INDEX_LOCK = threading.RLock()
TOKEN_RE = re.compile(r"[A-Za-z0-9_]{2,}")
_WHISPER_MODEL = None


def _now_iso():
    return datetime.now(timezone.utc).isoformat()


def _tokenize(text):
    return TOKEN_RE.findall((text or "").lower())


def _embed_text(text, dims=192):
    # Hashing trick embedding with L2 normalization.
    vec = [0.0] * dims
    for token in _tokenize(text):
        digest = hashlib.sha256(token.encode("utf-8")).digest()
        idx = int.from_bytes(digest[:4], "big") % dims
        sign = 1.0 if (digest[4] & 1) == 0 else -1.0
        vec[idx] += sign
    norm = math.sqrt(sum(v * v for v in vec))
    if norm > 0:
        vec = [v / norm for v in vec]
    return vec


def _dot(a, b):
    return sum(x * y for x, y in zip(a, b))


def _safe_read_text(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _safe_read_pdf(path):
    try:
        import fitz
    except ImportError:
        return ""

    try:
        parts = []
        with fitz.open(path) as doc:
            for page in doc:
                page_text = page.get_text("text") or ""
                if page_text:
                    parts.append(page_text)
        return "\n".join(parts)
    except Exception:
        return ""



def _safe_read_image(path):
    try:
        import cv2
    except ImportError:
        return ""
    try:
        import pytesseract
    except ImportError:
        return ""

    try:
        image = cv2.imread(path)
        if image is None:
            return ""
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        denoised = cv2.GaussianBlur(gray, (3, 3), 0)
        _, thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        text = pytesseract.image_to_string(thresh)
        return text
    except Exception:
        return ""



def _safe_read_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        return ""

    try:
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join([p for p in parts if p])
    except Exception:
        return ""


def _safe_read_docx(path):
    try:
        import docx
    except ImportError:
        return ""

    try:
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception:
        return ""


def _safe_read_video(path):
    global _WHISPER_MODEL

    try:
        import whisper
    except ImportError:
        return ""

    ffmpeg_exe = os.getenv("FFMPEG_PATH", "").strip() or shutil.which("ffmpeg")
    if not ffmpeg_exe:
        try:
            import imageio_ffmpeg
            ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
        except Exception:
            ffmpeg_exe = ""

    if not ffmpeg_exe:
        return ""

    tmp_wav = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            tmp_wav = tmp.name

        cmd = [
            ffmpeg_exe,
            "-y",
            "-i",
            path,
            "-vn",
            "-acodec",
            "pcm_s16le",
            "-ar",
            "16000",
            "-ac",
            "1",
            tmp_wav,
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True)
        if proc.returncode != 0:
            return ""

        if _WHISPER_MODEL is None:
            _WHISPER_MODEL = whisper.load_model("base")

        result = _WHISPER_MODEL.transcribe(tmp_wav)
        return (result.get("text") or "").strip()
    except Exception:
        return ""
    finally:
        if tmp_wav and os.path.exists(tmp_wav):
            os.remove(tmp_wav)


def _safe_read_general_doc(path):
    try:
        from unstructured.partition.auto import partition
    except ImportError:
        return ""

    try:
        elements = partition(filename=path)
        text = "\n".join(str(element).strip() for element in elements if str(element).strip())
        return text
    except Exception:
        return ""


def _detect_file_kind(path):
    ext = os.path.splitext(path)[1].lower()
    mime, _ = mimetypes.guess_type(path)
    mime = (mime or "").lower()

    if mime.startswith("image/") or ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".gif"}:
        return "image"
    if mime.startswith("video/") or ext in {".mp4", ".mov", ".avi", ".mkv", ".webm"}:
        return "video"
    if mime == "application/pdf" or ext == ".pdf":
        return "pdf"
    if mime.startswith("text/") or ext in {".txt", ".md", ".csv", ".json", ".py", ".js", ".html", ".css"}:
        return "text"
    return "general"


def parse_file(path):
    if not path or not os.path.exists(path):
        raise ValueError("File path does not exist.")

    ext = os.path.splitext(path)[1].lower()
    kind = _detect_file_kind(path)

    if kind == "text":
        return _safe_read_text(path)
    if kind == "pdf":
        text = _safe_read_pdf(path)
        if text:
            return text
        return _safe_read_general_doc(path)
    if kind == "image":
        return _safe_read_image(path)
    if kind == "video":
        return _safe_read_video(path)

    text = _safe_read_general_doc(path)
    if text:
        return text

    ext = os.path.splitext(path)[1].lower()
    if ext in {".docx", ".pptx", ".ppt", ".rtf", ".odt", ".doc"}:
        fallback = _safe_read_docx(path) if ext in {".docx"} else _safe_read_pptx(path)
        if fallback:
            return fallback

    raise ValueError(f"Unsupported or unreadable file type: {ext or 'unknown'}")

_parse_file = parse_file


def _chunk_text(text, chunk_size=900, overlap=120):
    clean = (text or "").strip()
    if not clean:
        return []

    chunk_size = max(300, int(chunk_size))
    overlap = max(0, min(int(overlap), chunk_size - 50))

    chunks = []
    start = 0
    length = len(clean)
    while start < length:
        end = min(start + chunk_size, length)
        if end < length:
            split_at = clean.rfind("\n", start, end)
            if split_at < start + 200:
                split_at = clean.rfind(" ", start, end)
            if split_at >= start + 200:
                end = split_at
        chunk = clean[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= length:
            break
        start = max(start + 1, end - overlap)
    return chunks


def _load_index(index_file):
    if not os.path.exists(index_file):
        return {"files": [], "chunks": []}
    with open(index_file, "r", encoding="utf-8") as f:
        return json.load(f)


def _save_index(index_file, data):
    os.makedirs(os.path.dirname(index_file) or ".", exist_ok=True)
    with open(index_file, "w", encoding="utf-8") as f:
        json.dump(data, f)


def index_file(path, index_file, vector_dims=192, chunk_size=900, chunk_overlap=120):
    filename = os.path.basename(path)
    text = _parse_file(path)
    chunks = _chunk_text(text, chunk_size=chunk_size, overlap=chunk_overlap)
    if not chunks:
        raise ValueError(f"No text could be extracted from {filename}.")

    file_id = str(uuid.uuid4())
    file_record = {
        "file_id": file_id,
        "name": filename,
        "path": path,
        "uploaded_at": _now_iso(),
        "size_bytes": os.path.getsize(path) if os.path.exists(path) else 0,
        "chunk_count": len(chunks),
        "text_excerpt": (text or "")[:1200],
    }

    with INDEX_LOCK:
        data = _load_index(index_file)
        data["files"].append(file_record)
        for idx, chunk in enumerate(chunks, start=1):
            data["chunks"].append({
                "chunk_id": f"{file_id}:{idx}",
                "file_id": file_id,
                "file_name": filename,
                "chunk_index": idx,
                "text": chunk,
                "vector": _embed_text(chunk, dims=vector_dims),
            })
        _save_index(index_file, data)

    return file_record


def search_index(query, index_file, top_k=4, vector_dims=192, file_ids=None):
    query = (query or "").strip()
    if not query:
        return []

    with INDEX_LOCK:
        data = _load_index(index_file)
        chunks = data.get("chunks", [])

    if not chunks:
        return []

    q_vec = _embed_text(query, dims=vector_dims)
    scored = []
    file_ids_set = set(file_ids or [])
    for chunk in chunks:
        if file_ids_set and chunk.get("file_id") not in file_ids_set:
            continue
        vec = chunk.get("vector") or []
        if len(vec) != vector_dims:
            continue
        score = _dot(q_vec, vec)
        if score > 0.05:
            scored.append({
                "score": score,
                "file_id": chunk["file_id"],
                "file_name": chunk["file_name"],
                "chunk_index": chunk["chunk_index"],
                "text": chunk["text"],
            })

    scored.sort(key=lambda item: item["score"], reverse=True)
    return scored[: max(1, int(top_k))]


def list_documents(index_file):
    with INDEX_LOCK:
        data = _load_index(index_file)
    files = data.get("files", [])
    return sorted(files, key=lambda item: item.get("uploaded_at", ""), reverse=True)
